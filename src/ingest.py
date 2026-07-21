# src/ingest.py

import os
import fitz
import pickle
import tiktoken
import shutil
from dotenv import load_dotenv
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from langchain_openai import AzureOpenAIEmbeddings
from langchain_postgres import PGVector
from rank_bm25 import BM25Okapi
from src.metrics import StageTimer, TokenUsage

load_dotenv()

CONNECTION_STRING = os.getenv("DATABASE_URL")
COLLECTION_NAME   = "documind_chunks"
MAX_PAGES         = int(os.getenv("MAX_PDF_PAGES", 20))

ALLOWED_EXTENSIONS = {".pdf", ".txt", ".md", ".docx", ".pptx"}

# ── Embeddings ────────────────────────────────────────────────────────────────

embeddings = AzureOpenAIEmbeddings(
    azure_deployment=os.getenv("AZURE_OPENAI_EMBEDDING_DEPLOYMENT"),
    azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
    api_key=os.getenv("AZURE_OPENAI_API_KEY"),
    api_version=os.getenv("AZURE_OPENAI_API_VERSION"),
)

# ── Splitter ──────────────────────────────────────────────────────────────────

splitter = RecursiveCharacterTextSplitter(
    chunk_size=512,
    chunk_overlap=100,
    separators=["\n\n", "\n", ".", " "]
)

# ── Extractors ────────────────────────────────────────────────────────────────

def extract_pdf(file_path: str, session_id: str) -> list[Document]:
    doc    = fitz.open(file_path)
    pages  = []
    for i, page in enumerate(doc):
        text = page.get_text().strip()
        if text:
            pages.append(Document(
                page_content=text,
                metadata={
                    "source":     os.path.basename(file_path),
                    "page":       i + 1,
                    "session_id": session_id
                }
            ))
    doc.close()
    return pages


def extract_txt(file_path: str, session_id: str) -> list[Document]:
    with open(file_path, "r", encoding="utf-8") as f:
        text = f.read().strip()
    return [Document(
        page_content=text,
        metadata={
            "source":     os.path.basename(file_path),
            "page":       1,
            "session_id": session_id
        }
    )]


def extract_docx(file_path: str, session_id: str) -> list[Document]:
    from docx import Document as DocxDocument
    doc   = DocxDocument(file_path)
    pages = []
    page_text = []
    page_num  = 1

    for para in doc.paragraphs:
        if para.text.strip():
            page_text.append(para.text.strip())
        # treat every 15 paragraphs as a "page" for metadata
        if len(page_text) >= 15:
            pages.append(Document(
                page_content="\n".join(page_text),
                metadata={
                    "source":     os.path.basename(file_path),
                    "page":       page_num,
                    "session_id": session_id
                }
            ))
            page_text = []
            page_num += 1

    # remaining paragraphs
    if page_text:
        pages.append(Document(
            page_content="\n".join(page_text),
            metadata={
                "source":     os.path.basename(file_path),
                "page":       page_num,
                "session_id": session_id
            }
        ))
    return pages


def extract_pptx(file_path: str, session_id: str) -> list[Document]:
    from pptx import Presentation
    prs   = Presentation(file_path)
    pages = []
    for i, slide in enumerate(prs.slides):
        text = "\n".join([
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ])
        if text:
            pages.append(Document(
                page_content=text,
                metadata={
                    "source":     os.path.basename(file_path),
                    "page":       i + 1,
                    "session_id": session_id
                }
            ))
    return pages


def extract_document(file_path: str, session_id: str) -> list[Document]:
    """Router — picks the right extractor based on file extension."""
    ext = os.path.splitext(file_path)[1].lower()
    extractors = {
        ".pdf":  extract_pdf,
        ".txt":  extract_txt,
        ".md":   extract_txt,   # same as txt
        ".docx": extract_docx,
        ".pptx": extract_pptx,
    }
    if ext not in extractors:
        raise ValueError(f"Unsupported file type: {ext}. Allowed: {', '.join(ALLOWED_EXTENSIONS)}")
    return extractors[ext](file_path, session_id)

# ── Chunking ──────────────────────────────────────────────────────────────────

def chunk_documents(pages: list[Document]) -> list[Document]:
    chunks = splitter.split_documents(pages)
    for i, chunk in enumerate(chunks):
        chunk.metadata["chunk_index"] = i
    print(f"  created {len(chunks)} chunks")
    return chunks

# ── BM25 ──────────────────────────────────────────────────────────────────────

def build_bm25_index(chunks: list[Document]) -> tuple:
    """Build BM25 index in memory — no disk persistence."""
    tokenized = [doc.page_content.lower().split() for doc in chunks]
    bm25      = BM25Okapi(tokenized)
    return bm25, chunks

# ── Vector store ──────────────────────────────────────────────────────────────

def store_in_pgvector(chunks: list[Document]) -> PGVector:
    from sqlalchemy import create_engine
    from sqlalchemy.pool import QueuePool

    engine = create_engine(
        CONNECTION_STRING,
        poolclass=QueuePool,
        pool_size=5,
        max_overflow=10,
        pool_pre_ping=True,
        pool_recycle=300,
    )
    vs = PGVector(
        embeddings=embeddings,
        collection_name=COLLECTION_NAME,
        connection=engine,
    )
    vs.add_documents(chunks, batch_size=len(chunks))
    return vs

# ── Main ingestion pipeline ───────────────────────────────────────────────────

def ingest_pdf(file_path: str, session_id: str) -> tuple:
    """
    Full ingestion pipeline.
    Accepts PDF, TXT, MD, DOCX, PPTX.
    Returns (vectorstore, bm25, docs)
    """
    filename = os.path.basename(file_path)
    ext      = os.path.splitext(filename)[1].lower()

    print(f"\nIngesting: {filename} [{ext}] for session {session_id[:8]}...")

    with StageTimer() as t:
        pages = extract_document(file_path, session_id)
    print(f"  extraction: {t.elapsed_ms:.0f}ms — {len(pages)} pages")

    # page limit check
    if len(pages) > MAX_PAGES:
        raise ValueError(
            f"Document has {len(pages)} pages. "
            f"Maximum allowed is {MAX_PAGES}. "
            f"Please upload a shorter document."
        )

    with StageTimer() as t:
        chunks = chunk_documents(pages)
    print(f"  chunking: {t.elapsed_ms:.0f}ms")

    # token count
    enc          = tiktoken.encoding_for_model("text-embedding-ada-002")
    exact_tokens = sum(len(enc.encode(c.page_content)) for c in chunks)
    usage        = TokenUsage(embedding_tokens=exact_tokens)
    print(f"  {exact_tokens} embedding tokens — ${usage.embedding_cost:.6f}")

    with StageTimer() as t:
        vs = store_in_pgvector(chunks)
    print(f"  embed+store: {t.elapsed_ms:.0f}ms")

    with StageTimer() as t:
        bm25, docs = build_bm25_index(chunks)
    print(f"  BM25 index: {t.elapsed_ms:.0f}ms")

    print(f"\nDone. {len(chunks)} chunks ready.\n")
    return vs, bm25, docs